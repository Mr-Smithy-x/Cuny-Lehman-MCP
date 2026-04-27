import math
import os
from pathlib import Path

from mcp.server.fastmcp import FastMCP
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Server Setup ───────────────────────────────────────────────────────────────
mcp = FastMCP("PowerPointProDesigner")

CWD: Path = Path(os.environ.get("FS_ROOT", os.getcwd())).resolve()
CWD.mkdir(parents=True, exist_ok=True)

# ── Design System ──────────────────────────────────────────────────────────────
# A curated set of named themes. Each has primary, secondary, accent,
# bg_dark, bg_light, text_dark, text_light colors.
THEMES = {
    "midnight": {
        "primary":     RGBColor(0x1E, 0x27, 0x61),
        "secondary":   RGBColor(0xCA, 0xDC, 0xFC),
        "accent":      RGBColor(0x4A, 0x90, 0xD9),
        "bg_dark":     RGBColor(0x12, 0x18, 0x3A),
        "bg_light":    RGBColor(0xF4, 0xF7, 0xFF),
        "text_dark":   RGBColor(0x1A, 0x1A, 0x2E),
        "text_light":  RGBColor(0xF0, 0xF4, 0xFF),
    },
    "forest": {
        "primary":     RGBColor(0x2C, 0x5F, 0x2D),
        "secondary":   RGBColor(0x97, 0xBC, 0x62),
        "accent":      RGBColor(0x4E, 0x9A, 0x51),
        "bg_dark":     RGBColor(0x1A, 0x38, 0x1B),
        "bg_light":    RGBColor(0xF5, 0xFA, 0xF0),
        "text_dark":   RGBColor(0x1A, 0x2E, 0x1B),
        "text_light":  RGBColor(0xF0, 0xFA, 0xF0),
    },
    "coral": {
        "primary":     RGBColor(0xF9, 0x61, 0x67),
        "secondary":   RGBColor(0xF9, 0xE7, 0x95),
        "accent":      RGBColor(0x2F, 0x3C, 0x7E),
        "bg_dark":     RGBColor(0x2F, 0x3C, 0x7E),
        "bg_light":    RGBColor(0xFF, 0xFA, 0xF8),
        "text_dark":   RGBColor(0x2F, 0x3C, 0x7E),
        "text_light":  RGBColor(0xFF, 0xFA, 0xF8),
    },
    "charcoal": {
        "primary":     RGBColor(0x36, 0x45, 0x4F),
        "secondary":   RGBColor(0xF2, 0xF2, 0xF2),
        "accent":      RGBColor(0x76, 0xB5, 0xC5),
        "bg_dark":     RGBColor(0x21, 0x21, 0x21),
        "bg_light":    RGBColor(0xFA, 0xFA, 0xFA),
        "text_dark":   RGBColor(0x21, 0x21, 0x21),
        "text_light":  RGBColor(0xF2, 0xF2, 0xF2),
    },
    "teal": {
        "primary":     RGBColor(0x02, 0x80, 0x90),
        "secondary":   RGBColor(0x00, 0xA8, 0x96),
        "accent":      RGBColor(0x02, 0xC3, 0x9A),
        "bg_dark":     RGBColor(0x01, 0x45, 0x55),
        "bg_light":    RGBColor(0xF0, 0xFA, 0xFB),
        "text_dark":   RGBColor(0x01, 0x2E, 0x36),
        "text_light":  RGBColor(0xF0, 0xFA, 0xFB),
    },
}
DEFAULT_THEME = "midnight"

# Slide canvas (standard widescreen 13.33" × 7.5")
SLIDE_W = 13.333
SLIDE_H = 7.5

# ── Resources ──────────────────────────────────────────────────────────────────
@mcp.resource(
    uri=CWD.as_uri(),
    name="CWD directory",
    description="The absolute path of the allowed filesystem.",
    mime_type="text/plain",
)
def cwd_resource() -> str:
    return str(CWD)

# ── Internal Helpers ───────────────────────────────────────────────────────────

def get_pptx(path: str) -> Presentation:
    """Load the existing presentation or create a new blank one."""
    return Presentation(path) if os.path.exists(path) else Presentation()


def resolve_theme(theme_name: str) -> dict:
    """Return theme dict by name, falling back to default."""
    return THEMES.get(theme_name, THEMES[DEFAULT_THEME])


def fit_font_size(
    lines: list[str],
    box_width_in: float,
    box_height_in: float,
    base_pt: int = 18,
    min_pt: int = 10,
    max_pt: int = 28,
    font_name: str = "Calibri",
) -> int:
    """
    Return the largest font size (pt) that fits all lines inside the given box.

    Uses a simple but reliable heuristic:
      - Estimate characters per line at each candidate pt size
        (Calibri is ~0.5× pt-width per char on average)
      - Sum wrapped lines × line height (1.35× pt for inter-line spacing)
      - Accept if total height ≤ box height
    """
    AVG_CHAR_WIDTH_RATIO = 0.50   # avg char width as fraction of pt size
    LINE_SPACING_RATIO   = 1.40   # line height = pt × this

    box_w_pts = box_width_in  * 72
    box_h_pts = box_height_in * 72

    for pt in range(max_pt, min_pt - 1, -1):
        char_width_pts  = pt * AVG_CHAR_WIDTH_RATIO
        chars_per_line  = max(1, int(box_w_pts / char_width_pts))
        line_height_pts = pt * LINE_SPACING_RATIO

        wrapped = sum(
            max(1, math.ceil(len(line) / chars_per_line)) if line.strip() else 1
            for line in lines
        )
        total_h = wrapped * line_height_pts
        if total_h <= box_h_pts:
            return pt

    return min_pt


def set_bg(slide, color: RGBColor):
    """Set a solid background color on a slide."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color: RGBColor):
    """Add a filled, borderless rectangle shape."""
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(
    slide,
    text: str,
    x: float, y: float, w: float, h: float,
    font_size: int,
    font_color: RGBColor,
    bold: bool = False,
    italic: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    font_name: str = "Calibri",
    word_wrap: bool = True,
) -> None:
    """Add a styled, non-overflowing text box to a slide."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap

    # Remove default interior margin so content aligns to box edges
    txBox.text_frame.margin_top    = Pt(0)
    txBox.text_frame.margin_bottom = Pt(0)
    txBox.text_frame.margin_left   = Pt(4)
    txBox.text_frame.margin_right  = Pt(4)

    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size    = Pt(font_size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = font_color
    run.font.name    = font_name


def add_bullet_textbox(
    slide,
    points: list[str],
    x: float, y: float, w: float, h: float,
    font_size: int,
    font_color: RGBColor,
    bullet_color: RGBColor,
    font_name: str = "Calibri",
) -> None:
    """Add a bullet-point text box that auto-sizes to fit the given box."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    txBox.text_frame.margin_top    = Pt(2)
    txBox.text_frame.margin_bottom = Pt(2)
    txBox.text_frame.margin_left   = Pt(4)
    txBox.text_frame.margin_right  = Pt(4)

    for i, point in enumerate(points):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        para.space_before = Pt(4)

        # Unicode bullet character
        run = para.add_run()
        run.text = f"\u2022  {point}"
        run.font.size  = Pt(font_size)
        run.font.color.rgb = font_color
        run.font.name  = font_name


# ── Creation Tools ─────────────────────────────────────────────────────────────

@mcp.tool()
def create_presentation(
    file_path: str,
    title: str,
    subtitle: str = "",
    theme: str = DEFAULT_THEME,
):
    """
    Creates a new PowerPoint with a professionally styled title slide.

    Args:
        file_path: Where to save the .pptx file.
        title:     Main title text.
        subtitle:  Optional subtitle / presenter name.
        theme:     Color theme name — one of: midnight, forest, coral, charcoal, teal.
    """
    t = resolve_theme(theme)
    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, t["bg_dark"])

    # Left accent strip
    add_rect(slide, 0, 0, 0.18, SLIDE_H, t["primary"])

    # Horizontal divider
    add_rect(slide, 0.38, SLIDE_H / 2 - 0.03, SLIDE_W - 0.76, 0.06, t["accent"])

    # Title — auto-sized to fit
    title_size = fit_font_size([title], box_width_in=11.5, box_height_in=1.4, max_pt=52, min_pt=28)
    add_textbox(
        slide, title,
        x=0.5, y=SLIDE_H / 2 - 1.7, w=12.3, h=1.5,
        font_size=title_size,
        font_color=t["text_light"],
        bold=True, align=PP_ALIGN.LEFT,
    )

    # Subtitle
    if subtitle:
        sub_size = fit_font_size([subtitle], box_width_in=11.5, box_height_in=0.6, max_pt=22, min_pt=12)
        add_textbox(
            slide, subtitle,
            x=0.5, y=SLIDE_H / 2 + 0.05, w=12.3, h=0.7,
            font_size=sub_size,
            font_color=t["secondary"],
            bold=False, align=PP_ALIGN.LEFT,
        )

    prs.save(file_path)
    return f"Created styled title slide → {file_path}"


@mcp.tool()
def add_section_title_slide(
    file_path: str,
    section_title: str,
    section_number: str = "",
    theme: str = DEFAULT_THEME,
):
    """
    Adds a visually distinct section-divider slide (dark background, large centered text).
    Use between major sections of a long deck.
    """
    t = resolve_theme(theme)
    prs = get_pptx(file_path)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, t["bg_dark"])

    if section_number:
        add_textbox(
            slide, section_number,
            x=1.0, y=2.2, w=11.3, h=0.7,
            font_size=20, font_color=t["secondary"],
            bold=False, align=PP_ALIGN.CENTER,
        )

    title_size = fit_font_size([section_title], box_width_in=11.0, box_height_in=1.4, max_pt=48, min_pt=24)
    add_textbox(
        slide, section_title,
        x=1.0, y=3.0, w=11.3, h=1.6,
        font_size=title_size, font_color=t["text_light"],
        bold=True, align=PP_ALIGN.CENTER,
    )

    # Thin centered underline accent
    add_rect(slide, 4.5, 4.75, 4.3, 0.05, t["accent"])

    prs.save(file_path)
    return f"Added section title slide to {file_path}"


@mcp.tool()
def add_bullet_slide(
    file_path: str,
    title: str,
    points: list[str],
    theme: str = DEFAULT_THEME,
):
    """
    Adds a clean bullet-point slide.
    Font size is automatically scaled so all bullets fit without overflowing.

    Args:
        file_path: Path to the .pptx file.
        title:     Slide title (36-40pt bold).
        points:    List of bullet point strings. Write as full sentences.
        theme:     Color theme name.
    """
    t = resolve_theme(theme)
    prs = get_pptx(file_path)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, t["bg_light"])

    # Left accent strip
    add_rect(slide, 0, 0, 0.12, SLIDE_H, t["primary"])

    # Title
    title_size = fit_font_size([title], box_width_in=11.8, box_height_in=1.0, max_pt=40, min_pt=22)
    add_textbox(
        slide, title,
        x=0.35, y=0.3, w=12.6, h=1.0,
        font_size=title_size, font_color=t["primary"],
        bold=True, align=PP_ALIGN.LEFT,
    )

    # Thin divider under title
    add_rect(slide, 0.35, 1.35, 12.6, 0.04, t["secondary"])

    # Bullets — fit_font_size accounts for number and length of points
    CONTENT_BOX_W = 12.4
    CONTENT_BOX_H = 5.6
    bullet_size = fit_font_size(points, box_width_in=CONTENT_BOX_W, box_height_in=CONTENT_BOX_H, max_pt=20, min_pt=10)
    add_bullet_textbox(
        slide, points,
        x=0.45, y=1.5, w=CONTENT_BOX_W, h=CONTENT_BOX_H,
        font_size=bullet_size, font_color=t["text_dark"], bullet_color=t["primary"],
    )

    prs.save(file_path)
    return f"Added bullet slide (font {bullet_size}pt) to {file_path}"


@mcp.tool()
def add_prose_slide(
    file_path: str,
    title: str,
    narrative_text: str,
    theme: str = DEFAULT_THEME,
):
    """
    Adds a slide with flowing narrative/prose text (no bullets).
    Ideal for executive summaries, context slides, or qualitative explanations.
    Font auto-sizes to prevent overflow.
    """
    t = resolve_theme(theme)
    prs = get_pptx(file_path)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, t["bg_light"])

    add_rect(slide, 0, 0, 0.12, SLIDE_H, t["primary"])

    title_size = fit_font_size([title], box_width_in=11.8, box_height_in=1.0, max_pt=40, min_pt=22)
    add_textbox(
        slide, title,
        x=0.35, y=0.3, w=12.6, h=1.0,
        font_size=title_size, font_color=t["primary"],
        bold=True, align=PP_ALIGN.LEFT,
    )
    add_rect(slide, 0.35, 1.35, 12.6, 0.04, t["secondary"])

    lines = narrative_text.split("\n")
    CONTENT_BOX_W = 12.4
    CONTENT_BOX_H = 5.6
    body_size = fit_font_size(lines, box_width_in=CONTENT_BOX_W, box_height_in=CONTENT_BOX_H, max_pt=18, min_pt=10)
    add_textbox(
        slide, narrative_text,
        x=0.45, y=1.5, w=CONTENT_BOX_W, h=CONTENT_BOX_H,
        font_size=body_size, font_color=t["text_dark"],
        bold=False, align=PP_ALIGN.LEFT,
        word_wrap=True,
    )

    prs.save(file_path)
    return f"Added prose slide (font {body_size}pt) to {file_path}"


@mcp.tool()
def add_two_column_slide(
    file_path: str,
    title: str,
    left_heading: str,
    left_points: list[str],
    right_heading: str,
    right_points: list[str],
    theme: str = DEFAULT_THEME,
):
    """
    Adds a two-column content slide — great for comparisons, pros/cons, before/after.
    Both columns auto-size their fonts independently to prevent overflow.
    """
    t = resolve_theme(theme)
    prs = get_pptx(file_path)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, t["bg_light"])

    add_rect(slide, 0, 0, 0.12, SLIDE_H, t["primary"])

    title_size = fit_font_size([title], box_width_in=12.5, box_height_in=0.9, max_pt=38, min_pt=20)
    add_textbox(
        slide, title,
        x=0.35, y=0.25, w=12.6, h=0.9,
        font_size=title_size, font_color=t["primary"],
        bold=True, align=PP_ALIGN.LEFT,
    )
    add_rect(slide, 0.35, 1.2, 12.6, 0.04, t["secondary"])

    COL_W = 5.9
    COL_H = 5.4
    COL_Y = 1.35
    LEFT_X  = 0.45
    RIGHT_X = 6.9

    # Column header backgrounds
    add_rect(slide, LEFT_X,  COL_Y, COL_W, 0.48, t["primary"])
    add_rect(slide, RIGHT_X, COL_Y, COL_W, 0.48, t["primary"])

    add_textbox(slide, left_heading,  LEFT_X,  COL_Y + 0.06, COL_W, 0.4, 16, t["text_light"], bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, right_heading, RIGHT_X, COL_Y + 0.06, COL_W, 0.4, 16, t["text_light"], bold=True, align=PP_ALIGN.CENTER)

    # Vertical separator
    add_rect(slide, SLIDE_W / 2 - 0.03, COL_Y, 0.06, COL_H, t["secondary"])

    BODY_Y = COL_Y + 0.6
    BODY_H = COL_H - 0.65

    left_size  = fit_font_size(left_points,  box_width_in=COL_W, box_height_in=BODY_H, max_pt=16, min_pt=9)
    right_size = fit_font_size(right_points, box_width_in=COL_W, box_height_in=BODY_H, max_pt=16, min_pt=9)

    add_bullet_textbox(slide, left_points,  LEFT_X,  BODY_Y, COL_W, BODY_H, left_size,  t["text_dark"], t["primary"])
    add_bullet_textbox(slide, right_points, RIGHT_X, BODY_Y, COL_W, BODY_H, right_size, t["text_dark"], t["primary"])

    prs.save(file_path)
    return f"Added two-column slide to {file_path}"


@mcp.tool()
def add_stat_callout_slide(
    file_path: str,
    title: str,
    stats: list[dict],
    theme: str = DEFAULT_THEME,
):
    """
    Adds a slide with large bold stat callouts — ideal for key metrics or KPI summaries.

    Args:
        file_path: Path to the .pptx file.
        title:     Slide title.
        stats:     List of dicts, each with keys:
                     "value"   — the big number/text (e.g. "42%", "$1.2M")
                     "label"   — descriptive label below the value
                   Max 4 stats per slide for legibility.
        theme:     Color theme name.

    Example stats:
        [{"value": "94%", "label": "Customer Satisfaction"},
         {"value": "$2.4M", "label": "Annual Revenue"},
         {"value": "3×", "label": "YoY Growth"}]
    """
    t = resolve_theme(theme)
    prs = get_pptx(file_path)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, t["bg_dark"])

    title_size = fit_font_size([title], box_width_in=12.5, box_height_in=1.0, max_pt=40, min_pt=22)
    add_textbox(slide, title, x=0.4, y=0.25, w=12.5, h=1.0, font_size=title_size,
                font_color=t["text_light"], bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 3.5, 1.3, 6.3, 0.05, t["accent"])

    stats = stats[:4]
    n = len(stats)
    card_w = (SLIDE_W - 1.2) / n - 0.3
    for i, stat in enumerate(stats):
        cx = 0.6 + i * (card_w + 0.3)
        cy = 1.8

        # Card background
        add_rect(slide, cx, cy, card_w, 4.2, t["primary"])

        val_size = fit_font_size([stat.get("value", "")], box_width_in=card_w - 0.2, box_height_in=1.6, max_pt=64, min_pt=28)
        add_textbox(slide, stat.get("value", ""),
                    cx + 0.1, cy + 0.5, card_w - 0.2, 1.7,
                    font_size=val_size, font_color=t["secondary"],
                    bold=True, align=PP_ALIGN.CENTER)

        add_rect(slide, cx + 0.4, cy + 2.4, card_w - 0.8, 0.04, t["accent"])

        lbl = stat.get("label", "")
        lbl_size = fit_font_size([lbl], box_width_in=card_w - 0.2, box_height_in=1.0, max_pt=16, min_pt=9)
        add_textbox(slide, lbl,
                    cx + 0.1, cy + 2.6, card_w - 0.2, 1.2,
                    font_size=lbl_size, font_color=t["text_light"],
                    bold=False, align=PP_ALIGN.CENTER)

    prs.save(file_path)
    return f"Added stat callout slide ({n} stats) to {file_path}"


@mcp.tool()
def add_chart_slide(
    file_path: str,
    title: str,
    categories: list[str],
    values: list[float],
    theme: str = DEFAULT_THEME,
):
    """Adds a styled Clustered Column Chart slide."""
    t = resolve_theme(theme)
    prs = get_pptx(file_path)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, t["bg_light"])

    add_rect(slide, 0, 0, 0.12, SLIDE_H, t["primary"])
    title_size = fit_font_size([title], box_width_in=12.5, box_height_in=1.0, max_pt=38, min_pt=20)
    add_textbox(slide, title, x=0.35, y=0.25, w=12.6, h=1.0,
                font_size=title_size, font_color=t["primary"], bold=True)
    add_rect(slide, 0.35, 1.3, 12.6, 0.04, t["secondary"])

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Data", values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5),
        chart_data
    ).chart

    # Style chart series color
    series = chart.series[0]
    fill = series.format.fill
    fill.solid()
    fill.fore_color.rgb = t["primary"]

    prs.save(file_path)
    return f"Added chart slide to {file_path}"


@mcp.tool()
def add_table_slide(
    file_path: str,
    title: str,
    rows: int,
    cols: int,
    data: list[list[str]],
    theme: str = DEFAULT_THEME,
):
    """
    Adds a slide with a styled table.
    Header row (row 0) gets a colored background with white text.
    """
    t = resolve_theme(theme)
    prs = get_pptx(file_path)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, t["bg_light"])

    add_rect(slide, 0, 0, 0.12, SLIDE_H, t["primary"])
    title_size = fit_font_size([title], box_width_in=12.5, box_height_in=1.0, max_pt=38, min_pt=20)
    add_textbox(slide, title, x=0.35, y=0.25, w=12.6, h=1.0,
                font_size=title_size, font_color=t["primary"], bold=True)
    add_rect(slide, 0.35, 1.3, 12.6, 0.04, t["secondary"])

    table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.6))
    tbl = table_shape.table

    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = str(data[r][c]) if r < len(data) and c < len(data[r]) else ""
            para = cell.text_frame.paragraphs[0]
            run = para.runs[0] if para.runs else para.add_run()
            run.font.size = Pt(13)
            run.font.bold = (r == 0)

            if r == 0:
                run.font.color.rgb = t["text_light"]
                cell.fill.solid()
                cell.fill.fore_color.rgb = t["primary"]
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = t["bg_light"]

    prs.save(file_path)
    return f"Added table slide to {file_path}"


@mcp.tool()
def add_image_slide(
    file_path: str,
    title: str,
    image_path: str,
    theme: str = DEFAULT_THEME,
):
    """Adds a slide with a title and a centered image."""
    t = resolve_theme(theme)
    prs = get_pptx(file_path)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, t["bg_light"])

    add_rect(slide, 0, 0, 0.12, SLIDE_H, t["primary"])
    title_size = fit_font_size([title], box_width_in=12.5, box_height_in=1.0, max_pt=38, min_pt=20)
    add_textbox(slide, title, x=0.35, y=0.25, w=12.6, h=1.0,
                font_size=title_size, font_color=t["primary"], bold=True)

    slide.shapes.add_picture(image_path, Inches(2.5), Inches(1.6), height=Inches(5.0))

    prs.save(file_path)
    return f"Added image slide to {file_path}"


@mcp.tool()
def add_quote_slide(
    file_path: str,
    quote: str,
    attribution: str = "",
    theme: str = DEFAULT_THEME,
):
    """
    Adds an impactful quote slide with large italic text on a dark background.
    Great for testimonials, key insights, or memorable statements.
    """
    t = resolve_theme(theme)
    prs = get_pptx(file_path)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, t["bg_dark"])

    # Decorative large opening quote mark
    add_textbox(slide, "\u201C", x=0.5, y=0.2, w=2.0, h=1.5,
                font_size=96, font_color=t["accent"], bold=False)

    quote_size = fit_font_size([quote], box_width_in=11.5, box_height_in=4.5, max_pt=32, min_pt=14)
    add_textbox(slide, quote,
                x=1.0, y=1.2, w=11.3, h=4.5,
                font_size=quote_size, font_color=t["text_light"],
                bold=False, italic=True, align=PP_ALIGN.CENTER, word_wrap=True)

    if attribution:
        attr_size = fit_font_size([attribution], box_width_in=11.3, box_height_in=0.6, max_pt=18, min_pt=10)
        add_textbox(slide, f"— {attribution}",
                    x=1.0, y=6.0, w=11.3, h=0.7,
                    font_size=attr_size, font_color=t["secondary"],
                    bold=False, align=PP_ALIGN.RIGHT)

    prs.save(file_path)
    return f"Added quote slide to {file_path}"


@mcp.tool()
def add_agenda_slide(
    file_path: str,
    items: list[str],
    theme: str = DEFAULT_THEME,
):
    """
    Adds a numbered agenda / table-of-contents slide.
    Each item gets a numbered circle badge and clear typography.
    """
    t = resolve_theme(theme)
    prs = get_pptx(file_path)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, t["bg_light"])

    add_rect(slide, 0, 0, 0.12, SLIDE_H, t["primary"])
    add_textbox(slide, "Agenda", x=0.35, y=0.25, w=12.6, h=0.9,
                font_size=38, font_color=t["primary"], bold=True)
    add_rect(slide, 0.35, 1.2, 12.6, 0.04, t["secondary"])

    items = items[:7]
    item_h = min(0.82, 5.8 / max(len(items), 1))
    item_size = fit_font_size(items, box_width_in=11.0, box_height_in=item_h * len(items), max_pt=18, min_pt=10)

    for i, item in enumerate(items):
        y = 1.35 + i * (item_h + 0.1)
        # Numbered circle (simulated with a small filled rectangle for portability)
        add_rect(slide, 0.4, y + 0.05, 0.52, item_h - 0.1, t["primary"])
        add_textbox(slide, str(i + 1), x=0.4, y=y + 0.05, w=0.52, h=item_h - 0.1,
                    font_size=14, font_color=t["text_light"], bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, item, x=1.05, y=y + 0.06, w=11.8, h=item_h,
                    font_size=item_size, font_color=t["text_dark"], bold=False)

    prs.save(file_path)
    return f"Added agenda slide ({len(items)} items) to {file_path}"


# ── Maintenance Tools ──────────────────────────────────────────────────────────

@mcp.tool()
def inspect_deck(file_path: str):
    """Returns a structural summary of the presentation (slide count, shape content)."""
    if not os.path.exists(file_path):
        return "File not found."
    prs = Presentation(file_path)
    out = []
    for i, s in enumerate(prs.slides):
        content = [shape.text if hasattr(shape, "text") else "[Object]" for shape in s.shapes]
        out.append(f"Slide {i}: {' | '.join(c for c in content if c.strip())[:120]}")
    return "\n".join(out)


@mcp.tool()
def read_slide_details(file_path: str, slide_idx: int):
    """Returns every shape on a slide with its index, type, and text content."""
    if not os.path.exists(file_path):
        return "File not found."
    prs = Presentation(file_path)
    if slide_idx >= len(prs.slides):
        return "Slide index out of range."
    slide = prs.slides[slide_idx]
    details = [f"Slide {slide_idx} details:"]
    for i, shape in enumerate(slide.shapes):
        text = shape.text if hasattr(shape, "text") else "[No text]"
        details.append(f"  Shape {i} | {shape.shape_type} | {text[:80]}")
    return "\n".join(details)


@mcp.tool()
def update_shape_text(file_path: str, slide_idx: int, shape_idx: int, new_text: str):
    """Replaces the text of an existing shape by index."""
    prs = get_pptx(file_path)
    try:
        shape = prs.slides[slide_idx].shapes[shape_idx]
        if not shape.has_text_frame:
            return "Error: shape has no text frame."
        shape.text_frame.text = new_text
        prs.save(file_path)
        return f"Updated slide {slide_idx}, shape {shape_idx}."
    except Exception as e:
        return f"Update failed: {e}"


@mcp.tool()
def delete_slide(file_path: str, slide_idx: int):
    """Removes a slide from the presentation by index."""
    prs = get_pptx(file_path)
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    if slide_idx >= len(slides):
        return "Error: slide index out of range."
    xml_slides.remove(slides[slide_idx])
    prs.save(file_path)
    return f"Deleted slide {slide_idx}."


@mcp.tool()
def clear_slide(file_path: str, slide_idx: int):
    """Removes all shapes from a slide but keeps the slide itself."""
    prs = get_pptx(file_path)
    slide = prs.slides[slide_idx]
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)
    prs.save(file_path)
    return f"Cleared all content from slide {slide_idx}."


@mcp.tool()
def list_themes() -> str:
    """Lists all available color theme names with a brief description."""
    descriptions = {
        "midnight": "Deep navy + ice blue — professional, corporate",
        "forest":   "Forest green + moss — nature, sustainability",
        "coral":    "Coral red + gold — energetic, bold",
        "charcoal": "Charcoal + off-white — minimal, elegant",
        "teal":     "Teal + seafoam — modern, tech, health",
    }
    return "\n".join(f"  {k}: {v}" for k, v in descriptions.items())


# ── Preview ────────────────────────────────────────────────────────────────────

@mcp.tool()
def export_preview_image(file_path: str, slide_idx: int, output_png: str):
    """Windows Only: Exports a single slide to PNG via PowerPoint COM automation."""
    try:
        #import comtypes.client
        #ppt = comtypes.client.CreateObject("PowerPoint.Application")
        #pres = ppt.Presentations.Open(os.path.abspath(file_path))
        #pres.Slides[slide_idx + 1].Export(os.path.abspath(output_png), "PNG")
        #pres.Close()
        return f"Preview failed (Windows/PowerPoint required)" #f"Preview saved to {output_png}"
    except Exception as e:
        return f"Preview failed (Windows/PowerPoint required): {e}"


# ── Entry Point ────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    mcp.run()